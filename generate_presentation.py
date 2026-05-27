from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

base = Path(__file__).resolve().parent
presentation = Presentation()

# Utility functions

def add_title_slide(title, subtitle):
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(title, lines, img_paths=None):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(24)
        p.font.name = 'Calibri'
    if img_paths:
        left = Inches(5.25)
        top = Inches(1.5)
        for img_path in img_paths:
            if img_path.exists():
                slide.shapes.add_picture(str(img_path), left, top, width=Inches(4))
                top += Inches(3)
            else:
                print(f'Missing image: {img_path}')
    return slide

def add_multi_image_slide(title, lines, img_paths):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22)
        p.font.name = 'Calibri'
    left = Inches(5)
    top = Inches(1.2)
    width = Inches(3.1)
    for img_path in img_paths:
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), left, top, width=width)
            top += Inches(2.4)
        else:
            print(f'Missing image: {img_path}')
    return slide

# Build slides
add_title_slide(
    'ResuMatch',
    'AI-powered resume matching platform for smarter hiring and job search'
)

add_content_slide(
    'Project Overview',
    [
        'ResuMatch is a complete job matching ecosystem for students, employers, and admins.',
        'It automatically generates resumes, matches jobs by skill similarity, and manages applications with role-based access.',
        'Designed for fast hiring decisions and better candidate-job fit.',
    ]
)

add_content_slide(
    'User Roles & Permissions',
    [
        '3 user types: Job Seekers, Employers, Admins.',
        'Admins have 5 levels with 50+ granular permissions.',
        'Employers post jobs, review applications, and accept or reject candidates.',
        'Job Seekers apply to jobs, see match scores, and download resumes.',
    ]
)

add_content_slide(
    'Core Features',
    [
        'Auto-generated resume from profile data.',
        'Multi-country resume formats with HTML preview and PDF export.',
        'Role-based access control and profile completion validation.',
        'Real-time notifications and dashboard refresh with Socket.io.',
        'Application tracking: pending, accepted, rejected.',
        'Pagination on lists for fast browsing and scalability.',
    ]
)

add_content_slide(
    'Problem Solved',
    [
        'Students often apply to unrelated jobs; ResuMatch gives better job recommendations.',
        'Employers receive more relevant candidates and spend less time filtering resumes.',
        'Admins can manage users, approvals, roles, and all platform data in one place.',
    ]
)

signup = base / 'frontend' / 'public' / 'signup.png'
login = base / 'frontend' / 'public' / 'login.png'
job1 = base / 'backend' / 'uploads' / 'job-images' / 'job-69bfec29a2ed83a6a8a23687-1775743106311-196729956.jpeg'
job2 = base / 'backend' / 'uploads' / 'job-images' / 'job-69bfec29a2ed83a6a8a23687-1775032577079-313825554.jpeg'
doc1 = base / 'backend' / 'uploads' / 'documents' / 'doc-69bc0887a292e50888de8305-1775113091058-970036790.jpeg'
profile1 = base / 'backend' / 'uploads' / 'profile-images' / 'profile-6a1596d496ee0096643dc5ab-1779809431288-622386418.jpeg'
add_multi_image_slide(
    'User Interface & Visuals',
    [
        'Clean signup/login screens for all users.',
        'Job listings with images and detailed descriptions.',
        'Profile and resume visuals show polished candidate experience.',
    ],
    [signup, login, job1, profile1]
)

add_content_slide(
    'How Matching Works',
    [
        'Employers enter job details and required skills during posting.',
        'Job Seekers fill out profiles and upload or auto-generate resumes.',
        'System computes TF-IDF vectors for resumes and job descriptions.',
        'Cosine similarity calculates a compatibility percentage score.',
        'Candidates are ranked and recommended by match strength.',
    ]
)

add_content_slide(
    'Resume System',
    [
        'Auto-generates resumes directly from user profile information.',
        'Supports 5 international resume formats for better market fit.',
        'Users can preview formats and convert to PDF via browser print.',
        'Ensures resume data stays consistent with the profile.',
    ]
)

employer_img = base / 'backend' / 'uploads' / 'job-images' / 'job-69bfec29a2ed83a6a8a23687-1775743761944-221431181.jpeg'
add_content_slide(
    'Employer Workflow',
    [
        'Employers create company profiles and publish jobs.',
        'They receive applications with a match score for each candidate.',
        'Employers can review resumes, accept, reject, and message applicants.',
        'Admin approval keeps job posts secure and accurate.',
    ],
    img_paths=[employer_img]
)

add_content_slide(
    'Admin & Security',
    [
        'Admins manage users, employers, jobs, and role assignments.',
        'Permission checks run before every action to protect data.',
        'Admin dashboard shows critical stats and approval workflows.',
        'Support for multi-level admin roles improves platform control.',
    ]
)

add_content_slide(
    'Tech Stack & Performance',
    [
        'Backend: Node.js, Express, MongoDB.',
        'Frontend: React with Vite for fast UI performance.',
        'Real-time Socket.io updates for notifications and dashboards.',
        'Indexes and pagination optimize search, listing, and dashboard speed.',
    ]
)

add_content_slide(
    'Business Benefits',
    [
        'Reduces time spent by employers reviewing irrelevant applications.',
        'Improves job search accuracy for candidates.',
        'Makes hiring decisions more data-driven and transparent.',
        'Supports both local and international resume needs.',
    ]
)

add_content_slide(
    'What’s Next',
    [
        'Native PDF resume export and richer resume templates.',
        'Email notifications and candidate messaging improvements.',
        'AI-driven resume optimization and interview preparation tools.',
        'Multi-language support and expanded analytics dashboard.',
    ]
)

add_content_slide(
    'Thank You',
    [
        'ResuMatch brings smarter matching, better resumes, and faster hiring.',
        'Ready to demo the platform and answer your questions.',
    ]
)

output = base / 'ResuMatch_Presentation.pptx'
presentation.save(output)
print('Created', output)
